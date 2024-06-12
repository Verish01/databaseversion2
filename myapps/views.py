from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, JsonResponse
from django import forms
from .models import Inventory
import csv
import os
import sqlite3
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import json
from django.utils.safestring import mark_safe
from .models import Weather
from .models import Movie

# Define InventoryForm class
class InventoryForm(forms.ModelForm):
    class Meta:
        model = Inventory
        fields = ['product_name', 'quantity_in_stock', 'cost_per_item', 'quantity_sold', 'sales', 'stock_date', 'photos']
        widgets = {
            'stock_date': forms.DateInput(attrs={'type': 'date'}),
        }

# Define QueryForm class
class QueryForm(forms.Form):
    query = forms.CharField(label='Enter your natural language query', max_length=200)

# Load the ML model and tokenizer
model_path = 'gaussalgo/T5-LM-Large-text2sql-spider'
model = AutoModelForSeq2SeqLM.from_pretrained(model_path)
tokenizer = AutoTokenizer.from_pretrained(model_path)

def query_database(db_path, query):
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        cursor.execute(query)
        results = cursor.fetchall()
        column_names = [description[0] for description in cursor.description]
        conn.close()
        return [dict(zip(column_names, row)) for row in results]
    except sqlite3.Error as e:
        return str(e)

def get_sql_query(question, schema):
    input_text = f"Question: {question} Schema: {schema}"
    model_inputs = tokenizer(input_text, return_tensors="pt")
    outputs = model.generate(**model_inputs, max_length=512)
    output_text = tokenizer.batch_decode(outputs, skip_special_tokens=True)
    return output_text[0]

def download_csv(results):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="results.csv"'
    writer = csv.writer(response)
    if results:
        headers = results[0].keys()
        writer.writerow(headers)
        for row in results:
            writer.writerow(row.values())
    return response

def download_word(results):
    document = Document()
    document.add_heading('Query Results', 0)
    if results:
        headers = results[0].keys()
        table = document.add_table(rows=1, cols=len(headers))
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(headers):
            hdr_cells[i].text = header
        for row in results:
            row_cells = table.add_row().cells
            for i, value in enumerate(row.values()):
                row_cells[i].text = str(value)
    document_path = 'query_results.docx'
    document.save(document_path)
    with open(document_path, 'rb') as docx_file:
        response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = 'attachment; filename=query_results.docx'
    return response

def download_ppt(results):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Results"
    if not results:
        return HttpResponse("No results found")
    headers = results[0].keys()
    table = slide.shapes.add_table(rows=len(results) + 1, cols=len(headers), left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(0.8)).table
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    for i, row in enumerate(results):
        for j, value in enumerate(row.values()):
            table.cell(i + 1, j).text = str(value)
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="results.pptx"'
    prs.save(response)
    return response

def clear_history(request):
    if 'query_history' in request.session:
        del request.session['query_history']
    return redirect('query_view')

def query_view(request):
    results = None
    sql_query = None
    query_history = request.session.get('query_history', [])
    if request.method == 'POST':
        form = QueryForm(request.POST)
        if form.is_valid():
            question = form.cleaned_data['query']
            # schema = '''"myapps_inventory" "DATABASE" 
            #             "product_ID" int, 
            #             "product_name" text, 
            #             "quantity_in_stock" int, 
            #             "cost_per_item" float, 
            #             primary key: "product_ID"'''
            
            # schema = '''"myapps_weather" "DATABASE" 
            #             "id" int, 
            #             "formatted_date" datetime, 
            #             "summary" varchar(200), 
            #             "precip_type" varchar(50), 
            #             "temperature_c" float, 
            #             "apparent_temperature_c" float, 
            #             "humidity" float, 
            #             "wind_speed_kmh" float, 
            #             "wind_bearing_degrees" int, 
            #             "visibility_km" float, 
            #             "cloud_cover" float, 
            #             "pressure_millibars" float, 
            #             "daily_summary" text, 
            #             primary key: "id"'''

            schema = '''"myapps_movie" "DATABASE" 
                        "id" int, 
                        "poster_link" varchar(200), 
                        "series_title" varchar(100), 
                        "released_year" varchar(10), 
                        "certificate" varchar(10), 
                        "runtime" varchar(20), 
                        "genre" varchar(100), 
                        "imdb_rating" float, 
                        "overview" text,
                        "meta_score" varchar(10), 
                        "director" varchar(100), 
                        "star1" varchar(100), 
                        "star2" varchar(100), 
                        "star3" varchar(100), 
                        "star4" varchar(100), 
                        "no_of_votes" bigint, 
                        "gross" varchar(50), 
                        primary key: "id"'''
            sql_query = get_sql_query(question, schema)
            db_path = os.path.join('db.sqlite3')
            results = query_database(db_path, sql_query)
            query_history.append({'question': question, 'sql_query': sql_query, 'results': results})
            request.session['query_history'] = query_history
            if 'csv' in request.POST:
                return download_csv(results)
            elif 'word' in request.POST:
                return download_word(results)
            elif 'ppt' in request.POST:
                return download_ppt(results)
            response_data = {
                'form': form, 
                'results': results, 
                'query': sql_query, 
                'query_history': query_history,
                'results_json': json.dumps(results)  # Pass the results as JSON
            }
            return render(request, 'query_view.html', response_data)
    else:
        form = QueryForm()
    return render(request, 'query_view.html', {'form': form, 'results': results, 'query': sql_query, 'query_history': query_history})

def inventory_list(request):
    inventory_items = Inventory.objects.all()
    return render(request, 'inventory_list.html', {'inventory_items': inventory_items})

def weather_list(request):
    weather_data = Weather.objects.all()
    return render(request, 'weather_list.html', {'weather_data': weather_data})

def weather_detail(request, pk):
    weather = get_object_or_404(Weather, pk=pk)
    return render(request, 'weather_detail.html', {'weather': weather})


def edit_item(request, id):
    item = get_object_or_404(Inventory, id=id)
    if request.method == "POST":
        form = InventoryForm(request.POST, instance=item)
        if form.is_valid():
            form.save()
            return redirect('inventory_list')
    else:
        form = InventoryForm(instance=item)
    return render(request, 'edit_item.html', {'form': form})

def delete_item(request, id):
    item = get_object_or_404(Inventory, id=id)
    if request.method == "POST":
        item.delete()
        return redirect('inventory_list')
    return render(request, 'confirm_delete.html', {'item': item})

def add_item(request):
    if request.method == 'POST':
        form = InventoryForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('inventory_list')
    else:
        form = InventoryForm()
    return render(request, 'add_item.html', {'form': form})

def inventory_visualization(request):
    return render(request, 'gradio.html')

def export_inventory_csv(request):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="inventory.csv"'
    writer = csv.writer(response)
    writer.writerow(['Product Name', 'Quantity in Stock', 'Cost per Item', 'Quantity Sold', 'Sales', 'Stock Date'])
    for item in Inventory.objects.all():
        writer.writerow([item.product_name, item.quantity_in_stock, item.cost_per_item, item.quantity_sold, item.sales, item.stock_date])
    return response

def export_weather_csv(request):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="weather.csv"'

    writer = csv.writer(response)
    writer.writerow(['Formatted Date', 'Summary', 'Precip Type', 'Temperature (C)', 
                     'Apparent Temperature (C)', 'Humidity', 'Wind Speed (km/h)', 
                     'Wind Bearing (degrees)', 'Visibility (km)', 'Cloud Cover', 
                     'Pressure (millibars)', 'Daily Summary'])

    for weather in Weather.objects.all():
        writer.writerow([
            weather.formatted_date, 
            weather.summary, 
            weather.precip_type, 
            weather.temperature_c, 
            weather.apparent_temperature_c, 
            weather.humidity, 
            weather.wind_speed_kmh, 
            weather.wind_bearing_degrees, 
            weather.visibility_km, 
            weather.cloud_cover, 
            weather.pressure_millibars, 
            weather.daily_summary
        ])

    return response

